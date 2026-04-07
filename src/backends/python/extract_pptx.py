#!/usr/bin/env python3
"""Extract per-slide content + tables from a .pptx file."""
import sys, json

def main():
    if len(sys.argv) < 2:
        print(json.dumps({"error": "Usage: extract_pptx.py <filepath>"}))
        sys.exit(1)
    filepath = sys.argv[1]
    try:
        from pptx import Presentation
    except ImportError:
        print(json.dumps({"error": "python-pptx not installed. Run: pip install python-pptx"}))
        sys.exit(1)
    try:
        prs = Presentation(filepath)
    except Exception as e:
        print(json.dumps({"error": f"Failed to open PPTX: {e}"}))
        sys.exit(1)
    slides = []
    for i, slide in enumerate(prs.slides):
        title = None
        text_parts = []
        tables = []
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                t = shape.text_frame.text.strip()
                if t:
                    is_title = False
                    try:
                        is_title = shape.is_placeholder and shape.placeholder_format.idx == 0
                    except (ValueError, AttributeError):
                        pass
                    if is_title:
                        title = t
                    else:
                        text_parts.append(t)
            if shape.has_table:
                rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                html = "<table>" + "".join("<tr>" + "".join(f"<td>{c}</td>" for c in row) + "</tr>" for row in rows) + "</table>"
                tables.append({"html": html})
        if not title and text_parts:
            title = text_parts.pop(0)
        notes = ""
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            pass
        full_text = "\n\n".join(filter(None, ([title] if title else []) + text_parts + ([f"Notes: {notes}"] if notes else [])))
        slides.append({"slide_idx": i, "title": title or "", "text": full_text, "tokens": len(full_text) // 4, "tables": tables})
    print(json.dumps({"slides": slides}))

if __name__ == "__main__":
    main()
